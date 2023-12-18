import pptx

def replace_keywords_in_pptx(pptx_path, keyword_dict, output_path):
    prs = pptx.Presentation(pptx_path)
    slovar = keyword_dict
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
        for k, v in slovar.items():
            if k in text_frame.text:
                text_frame.text = text_frame.text.replace(k,v)
    prs.save(output_path)

keyword_dict = {
    "${User_name}": "Кравченко Никита Глебович",
    "${Prof}": "Разработчик",
    "${%}": "95%"
}
input_file = "C:\\Users\\Администратор\\Desktop\\Сертификат.pptx"
output_file = "C:\\Users\\Администратор\\Desktop\\Сертификатoutput.pptx"
replace_keywords_in_pptx(input_file, keyword_dict, output_file)